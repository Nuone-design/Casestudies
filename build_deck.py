#!/usr/bin/env python3
"""
Case Study Deck v2 — matches dsr-deck-v16.html exactly.
17 slides, Manrope font, color-coded accent bars, two case studies.
1920×1080 canvas, Stripe/Linear type hierarchy.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Dimensions ─────────────────────────────────────────────
SLIDE_W_PX = 1920
SLIDE_H_PX = 1080
SLIDE_W = Emu(int(SLIDE_W_PX * 914400 / 96))  # px → EMU at 96 dpi
SLIDE_H = Emu(int(SLIDE_H_PX * 914400 / 96))

def px(v):
    """Convert pixel value to EMU (at 96 dpi)."""
    return Emu(int(v * 914400 / 96))

# ── Colors (from CSS :root) ────────────────────────────────
BG      = RGBColor(0x11, 0x11, 0x14)   # --s1
S2      = RGBColor(0x18, 0x18, 0x1C)   # --s2 surface/cards
S3      = RGBColor(0x1E, 0x1E, 0x23)   # --s3
B0      = RGBColor(0x1E, 0x1E, 0x22)   # ~rgba(255,255,255,0.06) on #111114
B1      = RGBColor(0x2C, 0x2C, 0x30)   # ~rgba(255,255,255,0.1)
TEXT    = RGBColor(0xF0, 0xED, 0xE8)   # --text
T2      = RGBColor(0xB8, 0xB5, 0xB0)   # --t2
T3      = RGBColor(0x88, 0x84, 0x80)   # --t3
T4      = RGBColor(0x55, 0x52, 0x50)   # --t4
BLUE    = RGBColor(0x4A, 0x8F, 0xD6)   # --blue
RED     = RGBColor(0xC4, 0x56, 0x2D)   # --red
GRN     = RGBColor(0x2A, 0x9E, 0x6E)   # --grn
AMB     = RGBColor(0xC0, 0x7A, 0x1A)   # --amb
PUR     = RGBColor(0x8A, 0x72, 0xD8)   # --pur
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

# Font
F = "Manrope"

# Act colors
ACT_COLORS = {
    "blue": BLUE, "red": RED, "grn": GRN, "amb": AMB, "pur": PUR
}

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helpers ────────────────────────────────────────────────

def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def accent_bar(slide, color_key):
    """5px left accent bar, full height."""
    c = ACT_COLORS.get(color_key, BLUE)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, px(0), px(0), px(5), SLIDE_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = c
    bar.line.fill.background()

def crumb(slide, text):
    """Breadcrumb text at top-left."""
    add_text(slide, 80, 56, 400, 30, text, size=18, weight=500, color=T4)

def label(slide, x, y, text):
    """Small tracked uppercase label."""
    tb = add_text(slide, x, y, 600, 28, text, size=18, weight=600, color=T4)
    for p in tb.text_frame.paragraphs:
        p.font.size = Pt(18)
    return tb

def add_text(slide, x, y, w, h, text, size=24, weight=400, color=TEXT,
             align=PP_ALIGN.LEFT, spacing=None, space_after=None):
    """Add text box with px coordinates."""
    txBox = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    txBox.text_frame.word_wrap = True
    tf = txBox.text_frame
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = weight >= 600
        p.font.name = F
        p.alignment = align
        if spacing:
            p.line_spacing = Pt(spacing)
        if space_after is not None:
            p.space_after = Pt(space_after)
    return txBox

def add_text_runs(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, spacing=None):
    """
    runs: list of (text, size, weight, color) tuples.
    All in one paragraph.
    """
    txBox = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    txBox.text_frame.word_wrap = True
    p = txBox.text_frame.paragraphs[0]
    p.alignment = align
    if spacing:
        p.line_spacing = Pt(spacing)
    for text, size, weight, color in runs:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = weight >= 600
        r.font.color.rgb = color
        r.font.name = F
    return txBox

def card(slide, x, y, w, h, fill=S2, border=B0, radius=0.06):
    """Rounded card."""
    c = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, px(x), px(y), px(w), px(h)
    )
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.line.color.rgb = border
    c.line.width = Pt(0.75)
    c.adjustments[0] = radius
    return c

def divider(slide, x, y, w, color=B0):
    """1px horizontal rule."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), Emu(9525)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

def image_zone(slide, x, y, w, h, hint="Drop image here"):
    """Dashed image placeholder zone."""
    iz = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, px(x), px(y), px(w), px(h)
    )
    iz.fill.solid()
    iz.fill.fore_color.rgb = S2
    iz.line.color.rgb = B1
    iz.line.width = Pt(1.5)
    iz.line.dash_style = 4  # Dash
    iz.adjustments[0] = 0.04
    # Hint text
    add_text(slide, x + w//2 - 120, y + h//2 - 20, 240, 40,
             hint, size=18, weight=500, color=T4, align=PP_ALIGN.CENTER)
    return iz

def nlist_item(slide, x, y, w, num, text, bold_parts=None):
    """Numbered list item with top border."""
    divider(slide, x, y, w)
    add_text(slide, x, y + 8, 40, 36, num, size=16, weight=700, color=T4)
    add_text(slide, x + 50, y + 6, w - 60, 50, text, size=24, weight=400, color=T2, spacing=38)
    return y + 60


# ══════════════════════════════════════════════════════════
# SLIDE 01 — ABOUT
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
accent_bar(s, "blue")

# Left side
label(s, 100, 140, "ABOUT")
add_text(s, 100, 190, 700, 70, "Nuwan", size=56, weight=700, color=TEXT)
add_text(s, 100, 280, 780, 200,
    "Lead/Staff Product Designer with a mixed in-house\n"
    "and agency background. I build with AI, not just\n"
    "design for it. A decade working on problems where\n"
    "design has to do more than look good.",
    size=24, weight=400, color=T2, spacing=40)

add_text_runs(s, 100, 500, 780, 40, [
    ("Host of ", 24, 400, T2),
    ("Designing Up", 24, 700, TEXT),
    (", a podcast on design leadership and AI in product teams.", 24, 400, T2),
])

# Pills
pills = ["B2B SaaS", "AI-native tools", "Design systems", "0\u21921 builds", "Platform products"]
pill_x = 100
for pl in pills:
    pw = len(pl) * 11 + 40
    pill = card(s, pill_x, 580, pw, 38, S2, B0, 0.5)
    add_text(s, pill_x + 10, 584, pw - 20, 30, pl, size=18, weight=500, color=T3, align=PP_ALIGN.CENTER)
    pill_x += pw + 10

# Right side — info cards
# Vertical split line
divider_v = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(960), px(0), Emu(9525), SLIDE_H)
divider_v.fill.solid()
divider_v.fill.fore_color.rgb = B0
divider_v.line.fill.background()

# Currently card
card(s, 1000, 140, 860, 140, S2, B0)
add_text(s, 1030, 155, 200, 22, "CURRENTLY", size=13, weight=700, color=T4)
add_text(s, 1030, 185, 800, 30, "Lead / Staff Product Designer", size=22, weight=700, color=TEXT)
add_text(s, 1030, 225, 800, 30, "Berlin, Germany", size=22, weight=400, color=T2)

# Background card
card(s, 1000, 300, 860, 180, S2, B0)
add_text(s, 1030, 315, 200, 22, "BACKGROUND", size=13, weight=700, color=T4)
add_text(s, 1030, 345, 800, 30, "In-house + Agency", size=22, weight=400, color=T2)
add_text(s, 1030, 380, 800, 30, "B2B SaaS + AI products", size=22, weight=400, color=T2)
add_text(s, 1030, 415, 800, 30, "Designing Up podcast", size=22, weight=400, color=T2)

# Outside work card
card(s, 1000, 500, 860, 140, S2, B0)
add_text(s, 1030, 515, 200, 22, "OUTSIDE WORK", size=13, weight=700, color=T4)
add_text(s, 1030, 545, 800, 30, "I run half marathons", size=22, weight=400, color=T2)
add_text(s, 1030, 580, 800, 30, "Addicted to Hyrox races", size=22, weight=400, color=T2)


# ══════════════════════════════════════════════════════════
# SLIDE 02 — LEADERSHIP
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
accent_bar(s, "pur")

label(s, 100, 80, "LEADERSHIP")
add_text(s, 100, 120, 900, 70, "How I operate", size=56, weight=700, color=TEXT)

# 3 columns
cols = [
    ("Set direction under ambiguity",
     "Defined clear principles and decision criteria to move teams forward before solutions were locked in."),
    ("Positioned design as a product partner",
     "Bridging executives, engineers and ICs around customer-centred decisions. A strategic voice, not a service function."),
    ("Scaled impact through people",
     "Established career frameworks and growth plans enabling clearer expectations, stronger ownership and team growth."),
]
for i, (title, body) in enumerate(cols):
    cx = 100 + i * 580
    # Top border
    divider(s, cx, 260, 520, B1)
    # Make the border slightly thicker
    add_text(s, cx, 288, 520, 60, title, size=26, weight=700, color=TEXT, spacing=34)
    add_text(s, cx, 370, 520, 160, body, size=20, weight=400, color=T2, spacing=33)

add_text(s, 100, 600, 900, 30,
    "Led the Product Design function at Spryker under the CPO",
    size=20, weight=400, color=T3)


# ══════════════════════════════════════════════════════════
# SLIDE 03 — DSR PROJECT COVER
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
accent_bar(s, "blue")

# Left
# Big faded number
add_text(s, 100, 180, 300, 100, "01", size=80, weight=800,
         color=RGBColor(0x22, 0x22, 0x26))

add_text(s, 100, 280, 300, 30, "SPRYKER", size=20, weight=700, color=T4)

add_text(s, 100, 330, 780, 120,
    "De-risking a \u20AC2M\nplatform bet",
    size=72, weight=800, color=TEXT, spacing=76)

add_text(s, 100, 510, 740, 200,
    "The company was under pressure to invest in a Digital\n"
    "Sales Room, driven by analyst hype, internal revenue\n"
    "narratives, and interest from a marquee enterprise\n"
    "customer \u2014 before there was evidence customers\n"
    "were ready to adopt it.",
    size=24, weight=400, color=T2, spacing=40)

# Vertical divider
dv = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(960), px(0), Emu(9525), SLIDE_H)
dv.fill.solid()
dv.fill.fore_color.rgb = B0
dv.line.fill.background()

# Right — browser frame placeholder
card(s, 1020, 120, 840, 840, RGBColor(0x1E, 0x1E, 0x24), B0)
# Browser bar
card(s, 1020, 120, 840, 52, RGBColor(0x28, 0x28, 0x2F), B0, 0.03)
# Traffic lights
for ci, clr in enumerate([RGBColor(0xFF, 0x5F, 0x57), RGBColor(0xFE, 0xBC, 0x2E), RGBColor(0x28, 0xC8, 0x40)]):
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, px(1042 + ci * 20), px(140), px(12), px(12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = clr
    dot.line.fill.background()

image_zone(s, 1030, 182, 820, 770, "Product screenshot")


# ══════════════════════════════════════════════════════════
# SLIDE 04 — MY ROLE
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
accent_bar(s, "blue")
crumb(s, "My role")

label(s, 100, 120, "WHAT I OWNED")
add_text(s, 100, 160, 700, 70, "My role", size=56, weight=700, color=TEXT)

items = [
    "Led end-to-end UX for the buyer-facing Digital Sales Room (DSR) product",
    "Sole designer \u2014 responsible for research, concepts, prototyping, stakeholder alignment",
    "Reported to CPO. Full design autonomy from first sketch to recommendation.",
]
iy = 280
for i, item in enumerate(items):
    iy = nlist_item(s, 100, iy, 780, f"0{i+1}", item)
    iy += 8

# Right — org cards
dv = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(960), px(0), Emu(9525), SLIDE_H)
dv.fill.solid()
dv.fill.fore_color.rgb = B0
dv.line.fill.background()

# My team card
card(s, 1000, 140, 860, 160, S2, B0)
add_text(s, 1030, 155, 200, 22, "MY TEAM", size=13, weight=700, color=T4)
add_text(s, 1030, 185, 800, 30, "Nuwan (Lead Designer)", size=22, weight=700, color=TEXT)
add_text(s, 1030, 220, 800, 30, "Solo practitioner on this workstream", size=22, weight=400, color=T2)

# Stakeholders card
card(s, 1000, 320, 860, 200, S2, B0)
add_text(s, 1030, 335, 200, 22, "STAKEHOLDERS", size=13, weight=700, color=T4)
add_text(s, 1030, 365, 800, 30, "CPO (direct report)", size=22, weight=400, color=T2)
add_text(s, 1030, 400, 800, 30, "Product Manager", size=22, weight=400, color=T2)
add_text(s, 1030, 435, 800, 30, "Engineering leads", size=22, weight=400, color=T2)
add_text(s, 1030, 470, 800, 30, "Sales leadership", size=22, weight=400, color=T2)

# Constraint card
card(s, 1000, 540, 860, 120, S2, B0)
add_text(s, 1030, 555, 200, 22, "CONSTRAINT", size=13, weight=700, color=T4)
add_text(s, 1030, 585, 800, 30, "De-risk without killing momentum", size=22, weight=400, color=T2)


# ══════════════════════════════════════════════════════════
# SLIDE 05 — THE PRESSURE (ACT 2)
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
accent_bar(s, "red")
crumb(s, "The pressure")

label(s, 100, 140, "ACT 02")
add_text(s, 100, 180, 900, 70, "The problems", size=56, weight=700, color=TEXT)

# 2x2 card grid
cards_data = [
    ("01", "Investor pressure", "New R&D investment needed. The market expected a move."),
    ("02", "Analyst hype", "Every B2B analyst was writing about digital sales rooms. FOMO was real."),
    ("03", "Sales enthusiasm", "Sales teams were pushing hard. They believed this was the edge they needed."),
    ("04", "Lighthouse customers", "Our best customers were interested. Some were aggressively asking us to build it."),
]
for i, (num, title, body) in enumerate(cards_data):
    col = i % 2
    row = i // 2
    cx = 100 + col * 440
    cy = 290 + row * 210
    card(s, cx, cy, 420, 190, S2, B0, 0.06)
    add_text(s, cx + 36, cy + 28, 100, 20, num, size=14, weight=700, color=T4)
    add_text(s, cx + 36, cy + 58, 360, 32, title, size=26, weight=700, color=TEXT)
    add_text(s, cx + 36, cy + 100, 360, 70, body, size=18, weight=400, color=T2, spacing=29)

# Bottom insight quote
divider(s, 100, 740, 860)
add_text(s, 100, 758, 860, 60,
    "Everyone had a reason to build. No one had a reason to pause. That was the problem.",
    size=22, weight=300, color=T3, spacing=37)


# ══════════════════════════════════════════════════════════
# SLIDE 06 — ASSUMPTIONS
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
accent_bar(s, "red")
crumb(s, "Starting assumptions")

# Left
label(s, 100, 120, "WHAT WE ASSUMED GOING IN")

stmts = [
    ("\u201CEnterprise complexity requires a\nshared digital workspace\u201D"),
    ("\u201CLighthouse customer interest signals\nbroader market demand\u201D"),
    ("\u201CA Digital Sales Room will shorten\ncycles and differentiate us\u201D"),
]
sy = 180
for i, stmt in enumerate(stmts):
    divider(s, 100, sy, 780)
    add_text(s, 100, sy + 10, 50, 40, f"0{i+1}", size=16, weight=700, color=T4)
    add_text(s, 160, sy + 8, 700, 80, stmt, size=36, weight=300, color=TEXT, spacing=47)
    sy += 120

add_text(s, 100, sy + 20, 780, 40,
    "These felt true. Several times, I almost convinced myself they were.",
    size=24, weight=400, color=T3)

# Right — image zone
dv = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(960), px(0), Emu(9525), SLIDE_H)
dv.fill.solid()
dv.fill.fore_color.rgb = B0
dv.line.fill.background()

image_zone(s, 1000, 100, 860, 880, "Hypothesis / research artifact")


# ══════════════════════════════════════════════════════════
# SLIDE 07 — DISCOVERY (ACT 3)
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
accent_bar(s, "grn")
crumb(s, "Discovery")

# Left
label(s, 100, 120, "THE WORK")
add_text(s, 100, 158, 780, 70, "Discovery and validation", size=56, weight=700, color=TEXT)

items = [
    "Synthesised all internal and external research and market data",
    "Profiled against ICP data to find genuine product-market matches",
    "Structured customer discovery calls across multiple verticals",
    "Mapped seller and buyer journeys separately. They were not the same problem.",
    "Daily synthesis shared with CPO. Full transparency throughout.",
]
iy = 280
for i, item in enumerate(items):
    iy = nlist_item(s, 100, iy, 780, f"0{i+1}", item)
    iy += 4

# Insight quote
divider(s, 100, iy + 10, 780)
add_text(s, 100, iy + 24, 780, 60,
    "Buyers were consistent. Sellers were chaos. Every customer ran different systems, cultures, processes.",
    size=22, weight=300, color=T3, spacing=37)

# Right — image zone
dv = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(960), px(0), Emu(9525), SLIDE_H)
dv.fill.solid()
dv.fill.fore_color.rgb = B0
dv.line.fill.background()

image_zone(s, 1000, 100, 860, 880, "Journey map / Miro board")


# ══════════════════════════════════════════════════════════
# SLIDE 08 — THE WORK (WHAT WE BUILT)
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
accent_bar(s, "grn")
crumb(s, "The work")

# Left
label(s, 100, 120, "CONCEPT AND TESTING")
add_text(s, 100, 158, 780, 70, "What we built and tested", size=56, weight=700, color=TEXT)

items = [
    "Designed a DSR \u2014 shared workspace where buyers and sellers collaborate on complex deals",
    "Tested key interactions: quote review, product configuration, video consultation booking",
    "Validated with 8 enterprise buyers across automotive and industrial verticals",
    "Rapid concept iteration. Each round tighter than the last.",
]
iy = 280
for i, item in enumerate(items):
    iy = nlist_item(s, 100, iy, 780, f"0{i+1}", item)
    iy += 4

divider(s, 100, iy + 10, 780)
add_text(s, 100, iy + 24, 780, 60,
    "The concept was strong. The problem was the market, not the design.",
    size=22, weight=300, color=T3, spacing=37)

# Right — image zone
dv = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(960), px(0), Emu(9525), SLIDE_H)
dv.fill.solid()
dv.fill.fore_color.rgb = B0
dv.line.fill.background()

image_zone(s, 1000, 100, 860, 880, "Wireframe / Hi-fi comp")


# ══════════════════════════════════════════════════════════
# SLIDE 09 — THE TURN (ACT 4)
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
accent_bar(s, "amb")
crumb(s, "The turn")

label(s, 100, 140, "WHAT THE RESEARCH ACTUALLY SAID")
add_text(s, 100, 180, 900, 70, "The assumptions were wrong", size=56, weight=700, color=TEXT)

# Reversal rows
reversals = [
    ("Enterprise complexity requires a shared workspace",
     "Adoption required transforming entrenched sales cultures overnight"),
    ("Lighthouse interest equals broader market demand",
     "One customer wanted a bespoke build. That is a consulting project, not a product."),
    ("A DSR will shorten sales cycles",
     "Customers had no CPQ, configurators or document management to even support it"),
]
ry = 310
for old, new in reversals:
    divider(s, 100, ry, 1720)
    # Old (struck through look — muted color)
    add_text(s, 100, ry + 12, 640, 60, old, size=22, weight=300, color=T4)
    # Arrow
    add_text(s, 760, ry + 14, 40, 30, "\u2192", size=20, weight=400, color=T4)
    # New (bold white)
    add_text(s, 820, ry + 12, 900, 60, new, size=22, weight=600, color=TEXT, spacing=33)
    ry += 90
divider(s, 100, ry, 1720)

# Big insight
divider(s, 100, ry + 40, 1720)
add_text(s, 100, ry + 60, 1200, 80,
    "Interest is not readiness.\nEnthusiasm is not adoption.",
    size=44, weight=800, color=TEXT, spacing=53)


# ══════════════════════════════════════════════════════════
# SLIDE 10 — THE RECOMMENDATION
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
accent_bar(s, "amb")
crumb(s, "The recommendation")

# Left
label(s, 100, 120, "THE DECISION")
add_text(s, 100, 158, 780, 70, "The recommendation", size=56, weight=700, color=TEXT)

items = [
    "Our ICP is not ready. Building ahead of readiness compounds the cost.",
    "Framed as business risk anchored to comparable project scope north of \u20AC2M",
    "Collaborated with infra PM to model integration, resource and roadmap costs",
    "CPO saw everything in real time. Recommendation landed without resistance.",
]
iy = 280
for i, item in enumerate(items):
    iy = nlist_item(s, 100, iy, 780, f"0{i+1}", item)
    iy += 4

divider(s, 100, iy + 10, 780)
add_text(s, 100, iy + 24, 780, 80,
    "I did not ask them to trust my instinct. I showed my working. Every claim was traceable back to a call, a concept, a data point.",
    size=22, weight=300, color=T3, spacing=37)

# Right — pull quote
dv = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(960), px(0), Emu(9525), SLIDE_H)
dv.fill.solid()
dv.fill.fore_color.rgb = B0
dv.line.fill.background()

add_text(s, 1040, 380, 200, 24, "CUSTOMER, VERBATIM", size=18, weight=500, color=T4)
add_text(s, 1040, 430, 780, 200,
    "\u201CGetting a tool like this adopted by our sales team off the bat will be tricky.\u201D",
    size=36, weight=300, color=T2, spacing=54)


# ══════════════════════════════════════════════════════════
# SLIDE 11 — DSR OUTCOME (ACT 5)
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
accent_bar(s, "pur")
crumb(s, "Outcome")

# Left
label(s, 100, 120, "WHAT HAPPENED")
add_text(s, 100, 158, 780, 70, "What happened", size=56, weight=700, color=TEXT)

metrics = [
    ("~\u20AC2M", "saved in avoided R&D investment"),
    ("8 mo", "of roadmap redirected to foundational capabilities"),
    ("Self-service shipped \u2713", "the redirect was the right bet"),
]
my = 280
for val, lbl in metrics:
    divider(s, 100, my, 780)
    sz = 60 if len(val) < 10 else 36
    add_text(s, 100, my + 10, 300, 65, val, size=sz, weight=800, color=TEXT)
    add_text(s, 420, my + 22, 460, 40, lbl, size=22, weight=400, color=T2)
    my += 85
divider(s, 100, my, 780)

add_text_runs(s, 100, my + 20, 780, 60, [
    ("Do differently: ", 24, 700, TEXT),
    ("Pressure-test the adoption assumption in week one. Ask the harder question before enthusiasm gets a budget.", 24, 400, T2),
])

# Right — closing quote
dv = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(960), px(0), Emu(9525), SLIDE_H)
dv.fill.solid()
dv.fill.fore_color.rgb = B0
dv.line.fill.background()

add_text(s, 1040, 340, 780, 300,
    "\u201CDesign\u2019s job here wasn\u2019t craft. It was making an invisible decision concrete enough that an entire organisation could act on it.\u201D",
    size=36, weight=300, color=T2, spacing=54)


# ══════════════════════════════════════════════════════════
# SLIDE 12 — DESIGN SYSTEM DIVIDER
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
accent_bar(s, "grn")

# Section divider — bottom-left aligned
add_text(s, 140, 680, 400, 30, "CASE STUDY 02", size=18, weight=700, color=T4)
add_text(s, 140, 720, 1000, 160,
    "Design System\nwith Claude Code",
    size=100, weight=800, color=TEXT, spacing=100)
add_text(s, 140, 920, 800, 60,
    "How I used 20% of my time and AI to fix what no one had time to fix,\nand created organisational pull in the process.",
    size=28, weight=400, color=T2, spacing=45)

# Giant faded number
add_text(s, 1500, 680, 400, 250, "02", size=240, weight=800,
         color=RGBColor(0x16, 0x16, 0x19))


# ══════════════════════════════════════════════════════════
# SLIDE 13 — DS HERO
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
accent_bar(s, "grn")

label(s, 160, 200, "CASE STUDY / DESIGN SYSTEM")
add_text(s, 160, 260, 1600, 400,
    "I 10x\u2019d myself\nwith AI to fix what\nno one had time\nto fix.",
    size=100, weight=800, color=TEXT, spacing=100)

# Green rule
rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(160), px(710), px(64), px(3))
rule.fill.solid()
rule.fill.fore_color.rgb = GRN
rule.line.fill.background()

add_text(s, 160, 750, 1200, 40,
    "Self-initiated. 20% of my time. One FE alongside me. Claude Code as the multiplier.",
    size=28, weight=400, color=T2)


# ══════════════════════════════════════════════════════════
# SLIDE 14 — DS PROBLEM
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
accent_bar(s, "grn")
crumb(s, "The problem")

# Left
label(s, 100, 120, "CONTEXT AND OWNERSHIP")
add_text(s, 100, 158, 780, 70, "The problem", size=56, weight=700, color=TEXT)

add_text(s, 100, 248, 780, 100,
    "Our demo shop is the first thing enterprise prospects see when evaluating a platform deal worth millions. The craft was poor. Inconsistent, patchy, nobody\u2019s job to fix. I couldn\u2019t accept that.",
    size=24, weight=400, color=T2, spacing=40)

items = [
    "Self-initiated. Pitched to my manager and got sales and demo team backing.",
    "20% of my time, one FE alongside me, Claude Code as the multiplier",
    "Three goals: elevate craft, increase prospect satisfaction, enable the team to ship faster without breaking",
]
iy = 400
for i, item in enumerate(items):
    iy = nlist_item(s, 100, iy, 780, f"0{i+1}", item)
    iy += 4

add_text(s, 100, iy + 16, 780, 30,
    "FULL VISIBILITY TO ALL STAKEHOLDERS THROUGHOUT",
    size=18, weight=600, color=GRN)

# Right — org cards
dv = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(960), px(0), Emu(9525), SLIDE_H)
dv.fill.solid()
dv.fill.fore_color.rgb = B0
dv.line.fill.background()

card(s, 1000, 140, 860, 140, S2, B0)
add_text(s, 1030, 155, 200, 22, "TEAM", size=13, weight=700, color=T4)
add_text(s, 1030, 185, 800, 30, "Lead Designer, Nuwan", size=22, weight=700, color=TEXT)
add_text(s, 1030, 220, 800, 30, "1 Frontend Engineer", size=22, weight=400, color=T2)

card(s, 1000, 300, 860, 220, S2, B0)
add_text(s, 1030, 315, 200, 22, "STAKEHOLDERS", size=13, weight=700, color=T4)
add_text(s, 1030, 345, 800, 30, "Manager", size=22, weight=400, color=T2)
add_text(s, 1030, 380, 800, 30, "Sales org", size=22, weight=400, color=T2)
add_text(s, 1030, 415, 800, 30, "Demo experience team", size=22, weight=400, color=T2)
add_text(s, 1030, 450, 800, 30, "Partner network", size=22, weight=400, color=T2)

card(s, 1000, 540, 860, 120, S2, B0)
add_text(s, 1030, 555, 200, 22, "CONSTRAINT", size=13, weight=700, color=T4)
add_text(s, 1030, 585, 800, 30, "Be ambitious. Don\u2019t break a live shop.", size=22, weight=400, color=T2)


# ══════════════════════════════════════════════════════════
# SLIDE 15 — DS HOW I BUILT IT
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
accent_bar(s, "grn")
crumb(s, "The work")

# Left
label(s, 100, 120, "APPROACH AND CRAFT")
add_text(s, 100, 158, 780, 70, "How I built it", size=56, weight=700, color=TEXT)

# 4 process steps
steps = [
    ("01", "Moodboard", "Visual direction and stakeholder alignment before a single component"),
    ("02", "Foundations", "Design principles, tokens, colour and type system"),
    ("03", "Components", "Built directly into codebase using Claude Code"),
    ("04", "Docs site", "Built simultaneously, not after"),
]
for i, (num, title, body) in enumerate(steps):
    sx = 100 + i * 195
    # Top border (green for first 3)
    border_color = GRN if i < 3 else B0
    bline = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(sx), px(280), px(175), px(3))
    bline.fill.solid()
    bline.fill.fore_color.rgb = border_color
    bline.line.fill.background()

    add_text(s, sx, 296, 175, 20, num, size=14, weight=700, color=T4)
    add_text(s, sx, 322, 175, 30, title, size=24, weight=700, color=TEXT)
    add_text(s, sx, 360, 175, 80, body, size=19, weight=400, color=T2, spacing=30)

divider(s, 100, 490, 780)
add_text(s, 100, 504, 780, 80,
    "The hardest decision was restraint. Every component had to feel like a step forward without breaking a live shop sales depends on daily.",
    size=22, weight=300, color=T3, spacing=37)

# Right — image zone
dv = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(960), px(0), Emu(9525), SLIDE_H)
dv.fill.solid()
dv.fill.fore_color.rgb = B0
dv.line.fill.background()

image_zone(s, 1000, 100, 860, 880, "Moodboard or tokens")


# ══════════════════════════════════════════════════════════
# SLIDE 16 — DS COMPONENTS
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
accent_bar(s, "grn")
crumb(s, "Components")

# Left
label(s, 100, 120, "THE OUTPUT")
add_text(s, 100, 158, 780, 70, "Components and documentation", size=56, weight=700, color=TEXT)

items = [
    "Foundation tokens applied system-wide across colour, type and spacing",
    "Component variations designed, tested and shared iteratively",
    "Documentation site built in parallel. A system nobody can navigate is just pretty Figma files.",
    "Claude Code handled the build layer. Taste, judgment and restraint stayed mine.",
]
iy = 300
for i, item in enumerate(items):
    iy = nlist_item(s, 100, iy, 780, f"0{i+1}", item)
    iy += 4

divider(s, 100, iy + 10, 780)
add_text(s, 100, iy + 24, 780, 80,
    "\u201CWhen AI does the building, the designer\u2019s job gets clearer. Taste, judgment, restraint. That\u2019s the 10% that matters.\u201D",
    size=22, weight=300, color=T3, spacing=37)

# Right — 2x image grid
dv = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(960), px(0), Emu(9525), SLIDE_H)
dv.fill.solid()
dv.fill.fore_color.rgb = B0
dv.line.fill.background()

image_zone(s, 1000, 100, 420, 430, "Component screen")
image_zone(s, 1440, 100, 420, 430, "Docs site screen")
image_zone(s, 1000, 550, 420, 430, "Token system")
image_zone(s, 1440, 550, 420, 430, "Storefront")


# ══════════════════════════════════════════════════════════
# SLIDE 17 — DS OUTCOME
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
accent_bar(s, "pur")
crumb(s, "Outcome")

# Left
label(s, 100, 120, "WHERE IT IS NOW")
add_text(s, 100, 158, 780, 70, "What happened", size=56, weight=700, color=TEXT)

metrics = [
    ("70%", "complete on 20% of my time"),
    ("0 breaks", "live shop stayed intact throughout"),
    ("Org started pulling \u2713", "sales, partners and cross-team allies appeared without me asking"),
]
my = 280
for val, lbl in metrics:
    divider(s, 100, my, 780)
    sz = 60 if len(val) < 15 else 36
    add_text(s, 100, my + 10, 300, 65, val, size=sz, weight=800, color=TEXT)
    add_text(s, 420, my + 22, 460, 40, lbl, size=22, weight=400, color=T2)
    my += 85
divider(s, 100, my, 780)

add_text_runs(s, 100, my + 20, 780, 60, [
    ("Do differently: ", 24, 700, TEXT),
    ("Start the documentation on day one. Undocumented early decisions become invisible debt in a system.", 24, 400, T2),
])

# Right — closing quote
dv = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(960), px(0), Emu(9525), SLIDE_H)
dv.fill.solid()
dv.fill.fore_color.rgb = B0
dv.line.fill.background()

add_text(s, 1040, 340, 780, 300,
    "\u201CI stopped having to push this forward. Sales started pushing my team to accelerate it. That is when I knew the signal was real.\u201D",
    size=36, weight=300, color=T2, spacing=54)


# ══════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════
output_path = "/Users/nuwanwithanage/Case study/Case_Study_Presentation.pptx"
prs.save(output_path)
print(f"\u2713 Saved to {output_path}")
print(f"  {len(prs.slides)} slides")
